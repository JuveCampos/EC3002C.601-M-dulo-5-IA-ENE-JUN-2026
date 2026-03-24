#!/usr/bin/env python3
"""
Genera la presentacion Clase07_ENOE.pptx
Estilo: 4:3, cintillas azul marino #1e4c7d, fuente Ubuntu, texto negro
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constantes de estilo ──────────────────────────────────────────────
AZUL = RGBColor(0x1E, 0x4C, 0x7D)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO = RGBColor(0x00, 0x00, 0x00)
GRIS = RGBColor(0x55, 0x55, 0x55)
VERDE = RGBColor(0x27, 0xAE, 0x60)
ROJO = RGBColor(0xC0, 0x39, 0x2B)
NARANJA = RGBColor(0xE6, 0x7E, 0x22)
MORADO = RGBColor(0x6C, 0x34, 0x83)
AMARILLO_OSCURO = RGBColor(0x7D, 0x66, 0x08)
BG_AZUL_CLARO = RGBColor(0xDD, 0xE5, 0xED)
BG_GRIS_CLARO = RGBColor(0xF0, 0xF0, 0xF0)
BG_HIGHLIGHT = RGBColor(0xEA, 0xF2, 0xFA)
BG_VERDE_CLARO = RGBColor(0xD5, 0xF5, 0xE3)
BG_ROJO_CLARO = RGBColor(0xFA, 0xDB, 0xD8)
BG_AMARILLO_CLARO = RGBColor(0xFE, 0xF9, 0xE7)
BG_MORADO_CLARO = RGBColor(0xF5, 0xEE, 0xF8)
BG_NARANJA_CLARO = RGBColor(0xFD, 0xEB, 0xD0)

FONT = "Ubuntu"
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Layout en blanco
blank_layout = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────
def add_blank_slide():
    return prs.slides.add_slide(blank_layout)


def add_banner(slide, text, x=0, y=0, w=None, h=Inches(0.65)):
    """Cintilla azul marino con texto blanco."""
    if w is None:
        w = Inches(5.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLANCO
    p.font.name = FONT
    return shape


def add_full_banner(slide, text, y=0, h=Inches(0.65)):
    return add_banner(slide, text, x=0, y=y, w=SLIDE_W, h=h)


def add_textbox(slide, text, x, y, w, h, font_size=18, bold=False,
                color=NEGRO, alignment=PP_ALIGN.LEFT, font_name=FONT,
                line_spacing=1.4):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rich_textbox(slide, x, y, w, h):
    """Devuelve el text_frame para agregar parrafos con formato mixto."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_run(paragraph, text, font_size=18, bold=False, italic=False,
            color=NEGRO, font_name=FONT):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return run


def new_paragraph(tf, font_size=18, bold=False, color=NEGRO,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=4):
    p = tf.add_paragraph()
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_colored_box(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rounded_box(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_table_slide(slide, data, col_widths, x, y, w, h):
    """data = lista de filas; primera fila = headers."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table

    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = FONT
                if row_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = BLANCO
                else:
                    p.font.color.rgb = NEGRO
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = AZUL
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    BLANCO if row_idx % 2 == 1
                    else RGBColor(0xF5, 0xF7, 0xFA)
                )
    return table_shape


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: PORTADA
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()

# Franja superior gris claro (logo area)
add_colored_box(s, 0, 0, SLIDE_W, Inches(1.2), BG_GRIS_CLARO)
add_textbox(s, "Tecnologico de Monterrey", Inches(0.5), Inches(0.25),
            Inches(4), Inches(0.8), font_size=22, bold=True, color=AZUL)

# Franja central azul claro
add_colored_box(s, 0, Inches(1.2), SLIDE_W, Inches(3.5), BG_AZUL_CLARO)
add_textbox(s, "07. ENOE", Inches(0.5), Inches(2.5), Inches(8), Inches(1.2),
            font_size=64, bold=True, color=AZUL)
add_textbox(s, "Modulo 5. Inteligencia Artificial", Inches(0.5), Inches(3.6),
            Inches(8), Inches(0.6), font_size=26, bold=True, color=AZUL)

# Zona inferior blanca con autor
tf = add_rich_textbox(s, Inches(6.5), Inches(5.3), Inches(3), Inches(1.8))
p = tf.paragraphs[0]
add_run(p, "Jorge Juvenal\nCampos Ferreira", font_size=22, bold=True)
p.alignment = PP_ALIGN.RIGHT
p2 = new_paragraph(tf, alignment=PP_ALIGN.RIGHT, space_before=8)
add_run(p2, "juvenal.campos@tec.mx", font_size=12, color=ROJO)

# Linea decorativa
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(6.3), Inches(5.3), Inches(0.04), Inches(1.5))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
line.line.fill.background()


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: QUE ES LA ENOE
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "ENOE", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, "La ")
add_run(p, "Encuesta Nacional de Ocupacion y Empleo (ENOE)", bold=True)
add_run(p, " es el principal instrumento estadistico que realiza el ")
add_run(p, "INEGI", bold=True)
add_run(p, " de forma continua (trimestral) en Mexico, con el objetivo de "
        "recopilar informacion detallada sobre las ")
add_run(p, "condiciones del mercado laboral", bold=True)
add_run(p, ": empleo, desempleo, informalidad y caracteristicas de la "
        "fuerza de trabajo.")

tf2 = add_rich_textbox(s, Inches(0.5), Inches(3.0), Inches(5), Inches(2.5))
p2 = tf2.paragraphs[0]
p2.line_spacing = Pt(24)
add_run(p2, "Su unidad de analisis son las ")
add_run(p2, "personas de 15 anos y mas", bold=True)
add_run(p2, " que residen en viviendas particulares "
        "seleccionadas en todo el territorio nacional.")
p3 = new_paragraph(tf2, space_before=12)
p3.line_spacing = Pt(24)
add_run(p3, "Inicio en el ")
add_run(p3, "primer trimestre de 2005", bold=True)
add_run(p3, ", reemplazando a la ENE (Encuesta Nacional de Empleo).")

# Info box
box = add_rounded_box(s, Inches(6), Inches(3.0), Inches(3.5), Inches(2.5),
                       BG_HIGHLIGHT)
tf3 = add_rich_textbox(s, Inches(6.3), Inches(3.3), Inches(3), Inches(2))
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "Empleo y Ocupacion\nen Mexico", font_size=16, bold=True,
        color=AZUL)
p2 = new_paragraph(tf3, alignment=PP_ALIGN.CENTER, space_before=12)
add_run(p2, "~120,000 viviendas\npor trimestre", font_size=13, color=GRIS)
p3 = new_paragraph(tf3, alignment=PP_ALIGN.CENTER, space_before=8)
add_run(p3, "~400,000 personas\nentrevistadas", font_size=13, color=GRIS)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: LETRAS E-N-O-E
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()

letters = [
    ("E", "Encuesta", RGBColor(0x24, 0x71, 0xA3)),
    ("N", "Nacional", RGBColor(0xC0, 0x39, 0x2B)),
    ("O", "Ocupacion", RGBColor(0x27, 0xAE, 0x60)),
    ("E", "Empleo", RGBColor(0xE6, 0x7E, 0x22)),
]

start_x = Inches(0.8)
for i, (letter, label, bg) in enumerate(letters):
    cx = start_x + Inches(i * 2.4)
    # Cuadro de color
    box = slide_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.8),
        Inches(1.6), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, letter, font_size=60, bold=True, color=BLANCO)

    # Label
    add_textbox(s, label, cx, Inches(3.55), Inches(1.6), Inches(0.5),
                font_size=20, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitulos informativos
info_items = [
    ("Indicadores de\nempleo y desempleo", Inches(1.0)),
    ("Publicacion\ntrimestral", Inches(3.8)),
    ("Seguimiento\npanel rotativo", Inches(6.6)),
]
for text, x in info_items:
    box = add_rounded_box(s, x, Inches(4.5), Inches(2.4), Inches(1.2),
                          BG_HIGHLIGHT)
    add_textbox(s, text, x + Inches(0.15), Inches(4.65),
                Inches(2.1), Inches(0.9),
                font_size=13, bold=True, alignment=PP_ALIGN.CENTER,
                color=AZUL)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: E - ENCUESTA
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "E - Encuesta", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, 'Una encuesta ("')
add_run(p, "E", bold=True)
add_run(p, '") es ')
add_run(p, "un ejercicio estadistico en el que se recopila informacion "
        "de una muestra", bold=True)
add_run(p, " \u2014un subconjunto de la poblacion\u2014 disenada para ser ")
add_run(p, "representativa", bold=True)
add_run(p, " del universo que se quiere estudiar.")

tf2 = add_rich_textbox(s, Inches(0.5), Inches(2.6), Inches(9), Inches(1.2))
p2 = tf2.paragraphs[0]
p2.line_spacing = Pt(26)
add_run(p2, "En el caso de la ENOE, INEGI selecciona aproximadamente ")
add_run(p2, "120,000 viviendas cada trimestre", bold=True)
add_run(p2, " distribuidas en todo el territorio nacional, y a partir "
        "de las respuestas de sus habitantes infiere las condiciones "
        "laborales de toda la poblacion de 15 anos y mas.")

# Box de ventajas
add_colored_box(s, Inches(0.5), Inches(4.2), Inches(0.08), Inches(2.8), AZUL)
box = add_rounded_box(s, Inches(0.6), Inches(4.2), Inches(8.8), Inches(2.8),
                       BG_HIGHLIGHT)
tf3 = add_rich_textbox(s, Inches(0.9), Inches(4.35), Inches(8.2), Inches(2.6))
p = tf3.paragraphs[0]
add_run(p, "Ventajas de la encuesta continua:", font_size=16, bold=True)
ventajas = [
    ("1. Frecuencia. ", "Se levanta cada trimestre, no cada 10 anos como un censo."),
    ("2. Profundidad tematica. ", "Se enfoca en el mercado laboral con gran detalle."),
    ("3. Oportunidad. ", "Resultados disponibles pocas semanas despues del cierre."),
    ("4. Seguimiento. ", "Permite observar cambios en el tiempo gracias al panel rotativo."),
]
for titulo, desc in ventajas:
    p = new_paragraph(tf3, space_before=4)
    add_run(p, titulo, font_size=15, bold=True)
    add_run(p, desc, font_size=15)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: N - NACIONAL
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "N - Nacional", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, 'La "')
add_run(p, "N", bold=True)
add_run(p, '" corresponde a ')
add_run(p, "Nacional", italic=True)
add_run(p, ". Esto enfatiza que l")
add_run(p, "os resultados de la ENOE tienen cobertura de todo el "
        "territorio mexicano y permiten realizar estimaciones a nivel pais",
        bold=True)
add_run(p, ".")

tf2 = add_rich_textbox(s, Inches(0.5), Inches(2.6), Inches(9), Inches(1.2))
p2 = tf2.paragraphs[0]
p2.line_spacing = Pt(26)
add_run(p2, "El diseno muestral tambien habilita ")
add_run(p2, "desagregaciones", bold=True, italic=True)
add_run(p2, " por ")
add_run(p2, "entidad federativa", bold=True)
add_run(p2, ", por ")
add_run(p2, "tamano de localidad", bold=True)
add_run(p2, ", por ")
add_run(p2, "ambitos rural/urbano", bold=True)
add_run(p2, " y por ")
add_run(p2, "39 ciudades autorrepresentadas", bold=True)
add_run(p2, ".")

# Columna izquierda
tf3 = add_rich_textbox(s, Inches(0.5), Inches(4.0), Inches(5), Inches(3))
p = tf3.paragraphs[0]
p.line_spacing = Pt(24)
add_run(p, "El caracter nacional le da ")
add_run(p, "legitimidad", bold=True, color=AZUL)
add_run(p, " para ser la ")
add_run(p, "fuente oficial", bold=True)
add_run(p, " de las estadisticas de empleo en Mexico.")

instituciones = [
    "Banxico \u2013 politica monetaria",
    "STPS \u2013 politica laboral",
    "CONEVAL \u2013 medicion de pobreza",
    "OIT \u2013 comparaciones internacionales",
    "Sector privado \u2013 analisis de mercado",
]
p = new_paragraph(tf3, space_before=10)
add_run(p, "Es utilizada por:", font_size=16)
for inst in instituciones:
    p = new_paragraph(tf3, space_before=4)
    add_run(p, "  \u2022 " + inst, font_size=15)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: O - OCUPACION
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "O - Ocupacion", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, 'La "')
add_run(p, "O", bold=True)
add_run(p, '" corresponde a ')
add_run(p, "Ocupacion", italic=True)
add_run(p, ". Este componente captura las ")
add_run(p, "actividades economicas", bold=True)
add_run(p, " que realizan las personas: que tipo de trabajo desempenan, "
        "en que sector, que posicion tienen, cuantas horas trabajan "
        "y cuanto ganan.")

tf2 = add_rich_textbox(s, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5))
p = tf2.paragraphs[0]
add_run(p, "La ENOE permite clasificar a la ")
add_run(p, "poblacion ocupada", bold=True)
add_run(p, " por:")

# Columna izquierda
items_left = [
    ("Sector de actividad", "Primario, secundario, terciario"),
    ("Posicion en la ocupacion", "Subordinados, cuenta propia, empleadores"),
    ("Tamano de establecimiento", "Micro, pequena, mediana, grande"),
]
items_right = [
    ("Nivel de ingresos", "En multiplos del salario minimo"),
    ("Formalidad / Informalidad", "Con o sin acceso a seguridad social"),
    ("Jornada laboral", "Horas trabajadas por semana"),
]

for col_items, x_off in [(items_left, Inches(0.5)), (items_right, Inches(5.0))]:
    y_off = Inches(3.5)
    for titulo, desc in col_items:
        tf = add_rich_textbox(s, x_off, y_off, Inches(4.3), Inches(1.0))
        p = tf.paragraphs[0]
        add_run(p, "\u2022 " + titulo, font_size=17, bold=True)
        p2 = new_paragraph(tf, space_before=2)
        add_run(p2, "   " + desc, font_size=14, color=GRIS)
        y_off += Inches(1.1)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: E - EMPLEO
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "E - Empleo", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.0))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, 'La segunda "')
add_run(p, "E", bold=True)
add_run(p, '" corresponde a ')
add_run(p, "Empleo", italic=True)
add_run(p, ". Clasifica a toda la poblacion de 15 anos y mas en "
        "categorias mutuamente excluyentes, siguiendo las "
        "recomendaciones de la ")
add_run(p, "OIT", bold=True)
add_run(p, ":")

# PEA Box
pea_box = add_rounded_box(s, Inches(0.4), Inches(2.5), Inches(4.4),
                           Inches(0.55), VERDE)
add_textbox(s, "Poblacion Economicamente Activa (PEA)",
            Inches(0.6), Inches(2.55), Inches(4), Inches(0.45),
            font_size=16, bold=True, color=BLANCO, alignment=PP_ALIGN.CENTER)

ocu_box = add_rounded_box(s, Inches(0.5), Inches(3.2), Inches(4.2),
                           Inches(0.6), BG_VERDE_CLARO)
add_textbox(s, "Ocupados: tienen un empleo\n(formal o informal)",
            Inches(0.7), Inches(3.25), Inches(3.8), Inches(0.5),
            font_size=14, bold=False)

des_box = add_rounded_box(s, Inches(0.5), Inches(3.95), Inches(4.2),
                           Inches(0.6), BG_ROJO_CLARO)
add_textbox(s, "Desocupados: no tienen empleo\npero lo buscan activamente",
            Inches(0.7), Inches(4.0), Inches(3.8), Inches(0.5),
            font_size=14, bold=False)

# PNEA Box
pnea_box = add_rounded_box(s, Inches(5.2), Inches(2.5), Inches(4.4),
                            Inches(0.55), RGBColor(0x7F, 0x8C, 0x8D))
add_textbox(s, "Poblacion No Economicamente Activa (PNEA)",
            Inches(5.3), Inches(2.55), Inches(4.2), Inches(0.45),
            font_size=14, bold=True, color=BLANCO, alignment=PP_ALIGN.CENTER)

disp_box = add_rounded_box(s, Inches(5.3), Inches(3.2), Inches(4.2),
                            Inches(0.6), RGBColor(0xEA, 0xEC, 0xEE))
add_textbox(s, "Disponibles: no buscan pero\naceptarian trabajar",
            Inches(5.5), Inches(3.25), Inches(3.8), Inches(0.5),
            font_size=14, bold=False)

nodisp_box = add_rounded_box(s, Inches(5.3), Inches(3.95), Inches(4.2),
                              Inches(0.6), RGBColor(0xF2, 0xF3, 0xF4))
add_textbox(s, "No disponibles: estudiantes,\njubilados, personas del hogar",
            Inches(5.5), Inches(4.0), Inches(3.8), Inches(0.5),
            font_size=14, bold=False)

# Nota OIT
add_textbox(s, "Esta clasificacion sigue las recomendaciones de la "
            "Organizacion Internacional del Trabajo (OIT).",
            Inches(0.5), Inches(5.0), Inches(9), Inches(0.4),
            font_size=13, color=GRIS, alignment=PP_ALIGN.CENTER)

# Diagrama conceptual abajo
add_colored_box(s, Inches(0.5), Inches(5.6), Inches(9), Inches(0.04), AZUL)

tf = add_rich_textbox(s, Inches(0.5), Inches(5.8), Inches(9), Inches(1.2))
p = tf.paragraphs[0]
add_run(p, "Indicadores derivados: ", font_size=15, bold=True, color=AZUL)
add_run(p, "Tasa de desocupacion (TD), Tasa de informalidad laboral (TIL), "
        "Tasa de subocupacion, Tasa de participacion, Tasa de condiciones "
        "criticas de ocupacion (TCCO), Tasa de presion general (TPRG).",
        font_size=14, color=GRIS)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: ENOE vs ENIGH
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "ENOE vs. ENIGH", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.6))
p = tf.paragraphs[0]
add_run(p, "Ambas son encuestas del INEGI con cobertura nacional, pero "
        "miden ")
add_run(p, "fenomenos distintos", bold=True)
add_run(p, ":")

data = [
    ["Aspecto", "ENOE", "ENIGH"],
    ["Objetivo", "Mercado laboral (empleo,\ndesempleo, informalidad)",
     "Ingresos y gastos\nde los hogares"],
    ["Frecuencia", "Trimestral (continua)", "Bienal (cada 2 anos)"],
    ["Unidad de analisis", "Personas de 15+ anos", "Hogares"],
    ["Muestra", "~120,000 viv./trimestre", "~90,000 viv./levantamiento"],
    ["Diseno temporal", "Panel rotativo (5 visitas)", "Corte transversal (1 visita)"],
    ["Ingresos", "Solo laborales\n(cuestionario corto)",
     "Todos los ingresos +\ngastos detallados"],
    ["Usuarios clave", "STPS, Banxico, OIT", "CONEVAL, politica social"],
]
add_table_slide(s, data,
                [Inches(2.0), Inches(3.5), Inches(3.5)],
                Inches(0.5), Inches(1.8), Inches(9), Inches(4.2))

# Resumen box
box = add_rounded_box(s, Inches(0.5), Inches(6.2), Inches(9), Inches(0.8),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(6.2), Inches(0.08), Inches(0.8), AZUL)
tf = add_rich_textbox(s, Inches(0.8), Inches(6.3), Inches(8.5), Inches(0.6))
p = tf.paragraphs[0]
add_run(p, "En resumen: ", font_size=15, bold=True)
add_run(p, "la ENOE mide ", font_size=15)
add_run(p, "quien trabaja, en que y como", font_size=15, italic=True)
add_run(p, "; la ENIGH mide ", font_size=15)
add_run(p, "cuanto ganan y en que gastan", font_size=15, italic=True)
add_run(p, " los hogares.", font_size=15)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: PANEL ROTATIVO
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Esquema de panel rotativo", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.0))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, "La ENOE utiliza un ")
add_run(p, "esquema de panel rotativo", bold=True)
add_run(p, " donde cada vivienda seleccionada es visitada durante ")
add_run(p, "5 trimestres consecutivos", bold=True)
add_run(p, ". Cada trimestre, se rota el ")
add_run(p, "20%", bold=True)
add_run(p, " de la muestra.")

# Diagrama de panel
colors_panel = [
    RGBColor(0x24, 0x71, 0xA3),
    RGBColor(0x27, 0xAE, 0x60),
    RGBColor(0xE6, 0x7E, 0x22),
    RGBColor(0x8E, 0x44, 0xAD),
    RGBColor(0xC0, 0x39, 0x2B),
    RGBColor(0x1A, 0xBC, 0x9C),
    RGBColor(0xF3, 0x9C, 0x12),
]
group_names = ["G-A", "G-B", "G-C", "G-D", "G-E", "G-F", "G-G"]
gris_panel = RGBColor(0xBD, 0xC3, 0xC7)

panel_data = [
    # (trimestre, [(grupo_idx, activo)])
    ("T1 2025", [(0, True), (1, True), (2, True), (3, True), (4, True)]),
    ("T2 2025", [(-1, False), (1, True), (2, True), (3, True), (4, True), (5, True)]),
    ("T3 2025", [(-1, False), (-1, False), (2, True), (3, True), (4, True), (5, True), (6, True)]),
]

bg_panel = add_rounded_box(s, Inches(0.4), Inches(2.5), Inches(9.2),
                            Inches(3.0), RGBColor(0xF8, 0xF9, 0xFA))

y_panel = Inches(2.7)
for trim, groups in panel_data:
    # Label
    add_textbox(s, trim, Inches(0.6), y_panel, Inches(1.2), Inches(0.45),
                font_size=13, bold=True, alignment=PP_ALIGN.RIGHT)

    x_cell = Inches(2.0)
    for g_idx, activo in groups:
        w_cell = Inches(1.05)
        h_cell = Inches(0.42)
        if activo and g_idx >= 0:
            box = add_rounded_box(s, x_cell, y_panel, w_cell, h_cell,
                                   colors_panel[g_idx])
            add_textbox(s, group_names[g_idx], x_cell, y_panel + Inches(0.02),
                        w_cell, h_cell - Inches(0.04),
                        font_size=12, bold=True, color=BLANCO,
                        alignment=PP_ALIGN.CENTER)
        else:
            box = add_rounded_box(s, x_cell, y_panel, w_cell, h_cell,
                                   gris_panel)
            add_textbox(s, "sale", x_cell, y_panel + Inches(0.02),
                        w_cell, h_cell - Inches(0.04),
                        font_size=11, color=RGBColor(0x99, 0x99, 0x99),
                        alignment=PP_ALIGN.CENTER)
        x_cell += Inches(1.15)

    # "nuevo" label for new groups
    if len(groups) > 5:
        add_textbox(s, "\u2190 nuevo", x_cell + Inches(0.05), y_panel,
                    Inches(0.8), Inches(0.4),
                    font_size=10, bold=True,
                    color=colors_panel[groups[-1][0]])

    y_panel += Inches(0.7)

# Explanation
add_textbox(s, "Cada grupo (G-A, G-B...) = ~20% de la muestra. En cualquier "
            "trimestre, 4 de 5 grupos se comparten con el trimestre anterior.",
            Inches(0.6), Inches(4.8), Inches(8.5), Inches(0.5),
            font_size=12, color=GRIS)

# Ventajas box
box = add_rounded_box(s, Inches(0.5), Inches(5.6), Inches(9), Inches(1.2),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(5.6), Inches(0.08), Inches(1.2), AZUL)
tf = add_rich_textbox(s, Inches(0.8), Inches(5.7), Inches(8.5), Inches(1.0))
p = tf.paragraphs[0]
add_run(p, "Ventajas del panel rotativo:", font_size=15, bold=True)
p2 = new_paragraph(tf, space_before=6)
add_run(p2, "1. Permite medir ", font_size=14)
add_run(p2, "transiciones laborales", font_size=14, italic=True)
add_run(p2, " (ej. de desocupado a ocupado).", font_size=14)
p3 = new_paragraph(tf, space_before=2)
add_run(p3, "2. Reduce la varianza de las estimaciones y distribuye "
        "la carga de trabajo de campo.", font_size=14)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: SEGUIMIENTO PANEL
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Seguimiento panel de trabajadores", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, "Gracias al panel rotativo, la ENOE permite hacer ")
add_run(p, "analisis de flujos laborales", bold=True)
add_run(p, ": observar como cambia la situacion de las ")
add_run(p, "mismas personas", bold=True)
add_run(p, " entre un trimestre y otro.")

# Flujos cards
flujos = [
    ("Desocupado \u2192 Ocupado", "Cuantos encontraron empleo",
     BG_VERDE_CLARO, RGBColor(0x1A, 0x7A, 0x3A)),
    ("Ocupado \u2192 Desocupado", "Cuantos perdieron su empleo",
     BG_ROJO_CLARO, ROJO),
    ("Informal \u2192 Formal", "Transicion a formalidad",
     BG_HIGHLIGHT, AZUL),
]

for i, (titulo, desc, bg, tc) in enumerate(flujos):
    cx = Inches(0.5) + Inches(i * 3.15)
    add_rounded_box(s, cx, Inches(2.4), Inches(2.9), Inches(1.5), bg)
    add_textbox(s, titulo, cx + Inches(0.2), Inches(2.6),
                Inches(2.5), Inches(0.5),
                font_size=15, bold=True, color=tc, alignment=PP_ALIGN.CENTER)
    add_textbox(s, desc, cx + Inches(0.2), Inches(3.2),
                Inches(2.5), Inches(0.4),
                font_size=13, color=GRIS, alignment=PP_ALIGN.CENTER)

tf2 = add_rich_textbox(s, Inches(0.5), Inches(4.3), Inches(9), Inches(0.8))
p = tf2.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, "Este seguimiento es posible porque las mismas viviendas son "
        "revisitadas hasta por ")
add_run(p, "5 trimestres (~15 meses)", bold=True)
add_run(p, ".")

box = add_rounded_box(s, Inches(0.5), Inches(5.3), Inches(9), Inches(1.0),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(5.3), Inches(0.08), Inches(1.0), AZUL)
tf3 = add_rich_textbox(s, Inches(0.8), Inches(5.4), Inches(8.5), Inches(0.8))
p = tf3.paragraphs[0]
add_run(p, "INEGI publica ", font_size=15)
add_run(p, "bases de datos de flujos laborales", font_size=15, bold=True)
add_run(p, " derivadas de este esquema panel, permitiendo estudiar ",
        font_size=15)
add_run(p, "trayectorias de empleo", font_size=15, italic=True)
add_run(p, " individuales.", font_size=15)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11: DISENO MUESTRAL
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Diseno muestral", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, "El diseno muestral de la ENOE es el conjunto de reglas y "
        "procedimientos que definen ")
add_run(p, "como se seleccionan las viviendas", bold=True)
add_run(p, " que seran entrevistadas.")

items_dm = [
    ("Probabilistico: ", "cada vivienda tiene una probabilidad conocida y "
     "distinta de cero de ser seleccionada."),
    ("Bietapico: ", "primero se seleccionan areas geograficas (UPM), "
     "luego viviendas dentro de ellas."),
    ("Estratificado: ", 'el pais se divide en "estratos" (por ejemplo, '
     "entidades federativas, urbano vs. rural) para asegurar que haya "
     "representacion de cada uno."),
    ("Por conglomerados: ", "se seleccionan primero areas geograficas "
     "(como manzanas o AGEBs), y dentro de ellas se eligen las viviendas."),
]

tf2 = add_rich_textbox(s, Inches(0.5), Inches(2.3), Inches(9), Inches(4.0))
for i, (titulo, desc) in enumerate(items_dm):
    p = tf2.paragraphs[0] if i == 0 else new_paragraph(tf2, space_before=10)
    p.line_spacing = Pt(24)
    add_run(p, "\u2022  ", font_size=18)
    add_run(p, titulo, font_size=18, bold=True)
    add_run(p, desc, font_size=18)

p_final = new_paragraph(tf2, space_before=18)
p_final.line_spacing = Pt(24)
add_run(p_final, "Este diseno determina tambien el tamano de la muestra "
        "(~120,000 viviendas por trimestre) y la precision estadistica "
        "de los resultados.", font_size=18)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12: FACTORES DE EXPANSION
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Factores de expansion", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
p = tf.paragraphs[0]
p.line_spacing = Pt(26)
add_run(p, "Cuando el INEGI selecciona una vivienda para la ")
add_run(p, "ENOE", bold=True)
add_run(p, ", esa vivienda no solo representa a si misma, sino a ")
add_run(p, "muchas otras viviendas similares en la poblacion", bold=True)
add_run(p, ". Para que los datos de la muestra \"se expandan\" al total "
        "de la poblacion, se asigna a cada persona un ")
add_run(p, "factor de expansion", bold=True)
add_run(p, ".")

items_fe = [
    "Imagina que se entrevisto a personas en 120,000 viviendas, pero en "
    "Mexico hay mas de 35 millones de viviendas.",
    'Cada persona de la muestra, entonces, debe "representar" a muchas '
    "personas reales.",
    "El factor de expansion indica cuantas personas de la poblacion "
    "equivale cada persona encuestada.",
]

tf2 = add_rich_textbox(s, Inches(0.5), Inches(2.8), Inches(9), Inches(2.5))
for i, item in enumerate(items_fe):
    p = tf2.paragraphs[0] if i == 0 else new_paragraph(tf2, space_before=10)
    p.line_spacing = Pt(26)
    add_run(p, "\u2022  " + item, font_size=18)

# Box de ajustes
box = add_rounded_box(s, Inches(0.5), Inches(5.3), Inches(9), Inches(1.5),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(5.3), Inches(0.08), Inches(1.5), AZUL)
tf3 = add_rich_textbox(s, Inches(0.8), Inches(5.4), Inches(8.5), Inches(1.3))
p = tf3.paragraphs[0]
add_run(p, "Los factores se ajustan por: ", font_size=16, bold=True)
add_run(p, "probabilidad de seleccion, no respuesta y proyecciones "
        "demograficas del CONAPO.", font_size=16)
p2 = new_paragraph(tf3, space_before=8)
add_run(p2, "Es indispensable usar estos factores al calcular cualquier "
        "indicador a nivel poblacional.", font_size=15, bold=True, color=ROJO)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13: VARIABLES CLAVE
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_full_banner(s, "Variables clave del diseno muestral", y=Inches(0.3))

add_textbox(s, "Para declarar el diseno de encuesta necesitas estas "
            "variables presentes en las tablas de la ENOE:",
            Inches(0.5), Inches(1.2), Inches(9), Inches(0.6), font_size=18)

data_vars = [
    ["Variable", "Rol", "Descripcion"],
    ["fac_tri", "Peso muestral (weights)",
     "Factor de expansion trimestral"],
    ["fac_men", "Peso muestral mensual",
     "Factor para estimaciones mensuales"],
    ["upm", "Conglomerado",
     "Unidad Primaria de Muestreo (cluster)"],
    ["est_d_tri", "Estrato",
     "Estrato de diseno muestral trimestral"],
]
add_table_slide(s, data_vars,
                [Inches(2), Inches(3), Inches(4)],
                Inches(0.5), Inches(2.0), Inches(9), Inches(2.8))

add_textbox(s, "Nota: El nombre exacto de las variables puede variar entre "
            "ediciones. Siempre verifica con el descriptor de la base "
            "de datos correspondiente.",
            Inches(0.5), Inches(5.2), Inches(9), Inches(0.8),
            font_size=14, color=GRIS)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14: INDICADORES PRINCIPALES
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Indicadores principales", y=Inches(0.3))

add_textbox(s, "La ENOE produce los indicadores estrategicos del mercado "
            "laboral mexicano:",
            Inches(0.5), Inches(1.1), Inches(9), Inches(0.5), font_size=18)

indicadores = [
    ("Tasa de desocupacion (TD)", "% de la PEA sin empleo que busca activamente.",
     BG_VERDE_CLARO, RGBColor(0x1A, 0x7A, 0x3A)),
    ("Tasa de informalidad laboral (TIL)", "% de ocupados en condiciones de informalidad.",
     BG_ROJO_CLARO, ROJO),
    ("Tasa de subocupacion", "% de ocupados que buscan trabajar mas horas.",
     BG_HIGHLIGHT, AZUL),
    ("Tasa de participacion", "% de la poblacion 15+ economicamente activa.",
     BG_AMARILLO_CLARO, AMARILLO_OSCURO),
    ("Tasa de condiciones criticas (TCCO)", "Ocupados en condiciones precarias de trabajo.",
     BG_MORADO_CLARO, MORADO),
    ("Tasa de presion general (TPRG)", "Desocupados + ocupados que buscan otro empleo.",
     BG_NARANJA_CLARO, RGBColor(0xA0, 0x40, 0x00)),
]

for i, (titulo, desc, bg, tc) in enumerate(indicadores):
    col = i % 2
    row = i // 2
    cx = Inches(0.4) + Inches(col * 4.8)
    cy = Inches(1.8) + Inches(row * 1.6)
    add_rounded_box(s, cx, cy, Inches(4.5), Inches(1.3), bg)
    add_textbox(s, titulo, cx + Inches(0.2), cy + Inches(0.15),
                Inches(4.1), Inches(0.5),
                font_size=15, bold=True, color=tc)
    add_textbox(s, desc, cx + Inches(0.2), cy + Inches(0.65),
                Inches(4.1), Inches(0.5),
                font_size=13, color=NEGRO)

add_textbox(s, "Estos indicadores se publican a nivel nacional, por entidad "
            "federativa y por ciudad autorrepresentada.",
            Inches(0.5), Inches(6.8), Inches(9), Inches(0.4),
            font_size=13, color=GRIS, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 15: COVID - ETOE y ENOEN
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "COVID-19: ETOE y ENOEN", y=Inches(0.3))

add_textbox(s, "La pandemia de COVID-19 obligo a modificar el "
            "levantamiento de la ENOE:",
            Inches(0.5), Inches(1.1), Inches(9), Inches(0.5), font_size=18)

# Timeline
timeline_items = [
    ("Marzo 2020",
     "Se suspende el levantamiento presencial de la ENOE "
     "por la contingencia sanitaria."),
    ("Abril - Junio 2020",
     "Se implementa la ETOE (Encuesta Telefonica de Ocupacion "
     "y Empleo) como medida emergente. Resultados no "
     "directamente comparables con la ENOE."),
    ("Julio 2020",
     "Inicia la ENOEN (ENOE Nueva Edicion): modo mixto "
     "presencial + telefonico. Mismo marco conceptual."),
    ("2022 en adelante",
     "Regreso paulatino al levantamiento predominantemente "
     "presencial, con modalidad mixta integrada."),
]

# Linea vertical
add_colored_box(s, Inches(1.5), Inches(1.8), Inches(0.05), Inches(4.5), AZUL)

for i, (fecha, texto) in enumerate(timeline_items):
    cy = Inches(1.9) + Inches(i * 1.1)
    # Circulo
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(1.38), cy + Inches(0.05),
                               Inches(0.3), Inches(0.3))
    circ.fill.solid()
    circ.fill.fore_color.rgb = AZUL
    circ.line.fill.background()
    # Fecha
    add_textbox(s, fecha, Inches(2.0), cy, Inches(2.5), Inches(0.35),
                font_size=15, bold=True, color=AZUL)
    # Texto
    add_textbox(s, texto, Inches(2.0), cy + Inches(0.35), Inches(7.2),
                Inches(0.7), font_size=14, color=NEGRO)

# Nota importante
box = add_rounded_box(s, Inches(0.5), Inches(6.3), Inches(9), Inches(0.8),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(6.3), Inches(0.08), Inches(0.8), AZUL)
tf = add_rich_textbox(s, Inches(0.8), Inches(6.4), Inches(8.5), Inches(0.6))
p = tf.paragraphs[0]
add_run(p, "Importante: ", font_size=14, bold=True)
add_run(p, "al comparar series de tiempo, tener cuidado con el periodo "
        "abril-junio 2020 (ETOE), cuyos datos no son estrictamente "
        "comparables con la ENOE tradicional.", font_size=14)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 16: PREPARACION DE DATOS
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_full_banner(s, "Preparacion de los datos de la ENOE", y=Inches(0.3))

add_textbox(s, "La ENOE viene en multiples tablas (archivos CSV). "
            "Las principales son:",
            Inches(0.5), Inches(1.1), Inches(9), Inches(0.5), font_size=18)

data_tablas = [
    ["Tabla", "Contenido"],
    ["SDEM", "Variables sociodemograficas (sexo, edad, escolaridad, "
     "estado civil)"],
    ["COE1", "Condicion de actividad, busqueda de empleo, disponibilidad"],
    ["COE2", "Caracteristicas del empleo: ocupacion, sector, horas, ingresos"],
    ["HOGAR", "Caracteristicas de la vivienda y el hogar"],
    ["VIVIENDA", "Datos de identificacion y control de la vivienda"],
]
add_table_slide(s, data_tablas,
                [Inches(2), Inches(7)],
                Inches(0.5), Inches(1.7), Inches(9), Inches(3.0))

# Bloque de codigo
code_box = add_rounded_box(s, Inches(0.5), Inches(4.9), Inches(9),
                            Inches(2.2), RGBColor(0x1E, 0x2A, 0x3A))
tf = add_rich_textbox(s, Inches(0.8), Inches(5.0), Inches(8.4), Inches(2.0))
code_lines = [
    "# Cargar las tablas principales de la ENOE",
    "library(readr)",
    "",
    'sdem     <- read_csv("SDEMT125.csv")',
    'coe1     <- read_csv("COE1T125.csv")',
    'coe2     <- read_csv("COE2T125.csv")',
    'hogar    <- read_csv("HOGARTRIM125.csv")',
    'vivienda <- read_csv("VIVTRIM125.csv")',
]
for i, line in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else new_paragraph(tf, space_before=0,
                                                       space_after=0)
    c = RGBColor(0x7F, 0x8C, 0x8D) if line.startswith("#") else \
        RGBColor(0xE0, 0xE0, 0xE0)
    add_run(p, line, font_size=12, color=c, font_name="Courier New")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 17: CLASIFICADORES
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Clasificadores y catalogos", y=Inches(0.3))

add_textbox(s, "La ENOE utiliza clasificadores estandarizados para "
            "codificar las respuestas:",
            Inches(0.5), Inches(1.1), Inches(9), Inches(0.5), font_size=18)

clasif = [
    ("SINCO", "Sistema Nacional de\nClasificacion de Ocupaciones",
     "Clasifica los tipos de trabajo\nque realizan las personas.", AZUL),
    ("SCIAN", "Sistema de Clasificacion\nIndustrial de America del Norte",
     "Clasifica las actividades\neconomicas / sectores.", VERDE),
    ("Catalogos\ngeograficos", "Entidades, municipios\ny localidades",
     "Identifican la ubicacion\ngeografica de las viviendas.", NARANJA),
]

for i, (nombre, desc, detalle, color) in enumerate(clasif):
    cx = Inches(0.4) + Inches(i * 3.2)
    # Borde
    border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 cx, Inches(2.0), Inches(2.9), Inches(4.0))
    border.fill.solid()
    border.fill.fore_color.rgb = BLANCO
    border.line.color.rgb = color
    border.line.width = Pt(2)

    # Nombre
    add_textbox(s, nombre, cx + Inches(0.15), Inches(2.3),
                Inches(2.6), Inches(0.8),
                font_size=20, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)

    # Descripcion
    add_textbox(s, desc, cx + Inches(0.15), Inches(3.2),
                Inches(2.6), Inches(1.0),
                font_size=14, alignment=PP_ALIGN.CENTER)

    # Detalle
    add_textbox(s, detalle, cx + Inches(0.15), Inches(4.4),
                Inches(2.6), Inches(0.8),
                font_size=12, color=GRIS, alignment=PP_ALIGN.CENTER)

add_textbox(s, "Estos catalogos son fundamentales para interpretar "
            "correctamente los codigos numericos en las bases de datos.",
            Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
            font_size=17)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 18: DOCUMENTOS IMPORTANTES
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Documentos importantes", y=Inches(0.3))

docs = [
    ("Pagina general:",
     "https://www.inegi.org.mx/programas/enoe/15ymas/"),
    ("Microdatos:",
     "Seccion de Microdatos en la pagina del programa ENOE"),
    ("Tabulados interactivos:",
     "https://www.inegi.org.mx/app/tabulados/interactivos/"),
    ("Documentacion conceptual:",
     "Cuestionarios, guias de llenado y diseno muestral disponibles "
     "en la pagina del programa."),
    ("Glosario de la ENOE:",
     "Definiciones oficiales de los conceptos utilizados."),
]

for i, (titulo, desc) in enumerate(docs):
    cy = Inches(1.4) + Inches(i * 1.05)
    # Bullet circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(0.7), cy + Inches(0.08),
                               Inches(0.18), Inches(0.18))
    circ.fill.solid()
    circ.fill.fore_color.rgb = AZUL
    circ.line.fill.background()

    tf = add_rich_textbox(s, Inches(1.1), cy, Inches(8.2), Inches(0.9))
    p = tf.paragraphs[0]
    add_run(p, titulo + " ", font_size=17, bold=True)
    # URLs en azul subrayado
    if desc.startswith("http"):
        run = add_run(p, desc, font_size=16, color=AZUL)
    else:
        add_run(p, desc, font_size=16)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 19: RESUMEN
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_full_banner(s, "Resumen", y=Inches(0.3), h=Inches(0.7))

resumen_items = [
    ("Encuesta continua", "Se levanta cada trimestre\ndesde 2005"),
    ("Cobertura nacional", "32 entidades +\n39 ciudades"),
    ("Mercado laboral", "Empleo, desempleo,\ninformalidad"),
    ("Panel rotativo", "5 visitas por vivienda,\nflujos laborales"),
]

for i, (titulo, desc) in enumerate(resumen_items):
    cx = Inches(0.3) + Inches(i * 2.45)
    add_rounded_box(s, cx, Inches(1.5), Inches(2.2), Inches(2.5),
                    RGBColor(0xF8, 0xF9, 0xFA))
    add_textbox(s, titulo, cx + Inches(0.1), Inches(2.3),
                Inches(2.0), Inches(0.5),
                font_size=16, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, desc, cx + Inches(0.1), Inches(2.8),
                Inches(2.0), Inches(0.8),
                font_size=13, color=GRIS, alignment=PP_ALIGN.CENTER)

# Final box
add_rounded_box(s, Inches(0.5), Inches(4.5), Inches(9), Inches(1.2),
                BG_HIGHLIGHT)
tf = add_rich_textbox(s, Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.8))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "La ENOE es la fuente oficial para entender\nel mercado "
        "laboral de Mexico.", font_size=20, bold=True, color=AZUL)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 20: ERROR ESTANDAR E INCERTIDUMBRE
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Error estandar e incertidumbre", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.1), Inches(9), Inches(1.0))
p = tf.paragraphs[0]
p.line_spacing = Pt(24)
add_run(p, "Cuando estimamos un indicador a partir de una encuesta, "
        "el valor obtenido ")
add_run(p, "no es el valor exacto", bold=True)
add_run(p, " de la poblacion, sino una ")
add_run(p, "estimacion", bold=True)
add_run(p, " basada en una muestra. El ")
add_run(p, "error estandar", bold=True, color=RGBColor(0x1A, 0x7A, 0x3A))
add_run(p, " mide ")
add_run(p, "que tan dispersas estarian esas estimaciones", bold=True)
add_run(p, " si repitieramos el muestreo muchas veces.")

# Boxes comparativos
add_rounded_box(s, Inches(0.4), Inches(2.3), Inches(4.3), Inches(1.1),
                BG_VERDE_CLARO)
tf2 = add_rich_textbox(s, Inches(0.6), Inches(2.4), Inches(4), Inches(0.9))
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "Error estandar pequeno", font_size=15, bold=True,
        color=RGBColor(0x1A, 0x7A, 0x3A))
p2 = new_paragraph(tf2, alignment=PP_ALIGN.CENTER, space_before=4)
add_run(p2, "Estimacion estable y confiable.\n"
        "Cualquier muestra daria resultado similar.", font_size=13)

add_rounded_box(s, Inches(5.2), Inches(2.3), Inches(4.3), Inches(1.1),
                BG_ROJO_CLARO)
tf3 = add_rich_textbox(s, Inches(5.4), Inches(2.4), Inches(4), Inches(0.9))
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "Error estandar grande", font_size=15, bold=True, color=ROJO)
p2 = new_paragraph(tf3, alignment=PP_ALIGN.CENTER, space_before=4)
add_run(p2, "Alta variabilidad. Otra muestra podria\n"
        "dar un resultado muy diferente.", font_size=13)

# Box de importancia
box = add_rounded_box(s, Inches(0.5), Inches(3.7), Inches(9), Inches(3.2),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(3.7), Inches(0.08), Inches(3.2), AZUL)
tf4 = add_rich_textbox(s, Inches(0.8), Inches(3.8), Inches(8.5), Inches(3.0))
p = tf4.paragraphs[0]
add_run(p, "Por que es IMPORTANTE reportar la incertidumbre?",
        font_size=16, bold=True, color=AZUL)

razones = [
    ("1. Precision: ", "un resultado con error de 0.3% es mucho mas "
     "confiable que uno con error de 5%."),
    ("2. Intervalos de confianza: ", "el rango donde probablemente esta "
     "el valor real (ej. 55.6% \u00b1 1.2%)."),
    ("3. Comparaciones: ", "permite saber si la diferencia entre dos "
     "grupos es significativa o solo ruido muestral."),
    ("4. Estandar internacional: ", "OIT, Banco Mundial y buenas "
     "practicas exigen reportar variabilidad."),
]
for titulo, desc in razones:
    p = new_paragraph(tf4, space_before=6)
    add_run(p, titulo, font_size=15, bold=True)
    add_run(p, desc, font_size=15)

p_nota = new_paragraph(tf4, space_before=10)
add_run(p_nota, "En R, las funciones de srvyr (survey_mean, survey_total) "
        "calculan automaticamente el error estandar al declarar el "
        "diseno de encuesta.", font_size=13, italic=True, color=GRIS)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 21: IMPUTACION DE NAs
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Imputacion de valores faltantes (NAs)", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.8))
p = tf.paragraphs[0]
p.line_spacing = Pt(24)
add_run(p, "La ")
add_run(p, "imputacion", bold=True)
add_run(p, " es el proceso de reemplazar valores faltantes (NA) con "
        "valores estimados. En encuestas como la ENOE, los NAs surgen "
        "porque el informante no respondio, hubo errores de captura, "
        "o la pregunta no aplico.")

# Box de por que no eliminar
box = add_rounded_box(s, Inches(0.5), Inches(2.1), Inches(9), Inches(0.9),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(2.1), Inches(0.08), Inches(0.9), AZUL)
tf2 = add_rich_textbox(s, Inches(0.8), Inches(2.2), Inches(8.5), Inches(0.7))
p = tf2.paragraphs[0]
add_run(p, "Por que no simplemente eliminar los NAs? ", font_size=14,
        bold=True)
add_run(p, "Porque puede reducir la muestra y perder precision, e "
        "introducir sesgo si los datos faltantes no son aleatorios "
        "(ej. personas con ingresos altos tienden a no reportarlos).",
        font_size=14)

# Tabla de metodos
add_textbox(s, "Metodos comunes de imputacion:",
            Inches(0.5), Inches(3.2), Inches(9), Inches(0.4),
            font_size=16, bold=True)

data_imp = [
    ["Metodo", "Descripcion", "Ventaja", "Desventaja"],
    ["Media / Mediana", "Reemplaza NA con media\no mediana del grupo",
     "Simple, rapido", "Reduce variabilidad"],
    ["Regresion", "Predice el valor con\nun modelo", "Usa relaciones\n"
     "entre variables", "Puede sobreajustar"],
    ["KNN", "Usa los K vecinos\nmas cercanos", "Respeta relaciones\n"
     "locales", "Costoso en computo"],
    ["Hot Deck", "Reemplaza NA con valor\nde un donante similar",
     "Respeta distribucion\nreal", "Depende del matching"],
    ["Imputacion\nmultiple", "Genera M versiones\ny combina resultados",
     "Captura\nincertidumbre", "Mas complejo"],
]
add_table_slide(s, data_imp,
                [Inches(1.8), Inches(2.5), Inches(2.3), Inches(2.4)],
                Inches(0.5), Inches(3.7), Inches(9), Inches(3.5))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 22: METODO HOT DECK
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_banner(s, "Metodo Hot Deck", y=Inches(0.3))

tf = add_rich_textbox(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.7))
p = tf.paragraphs[0]
p.line_spacing = Pt(24)
add_run(p, "El metodo ")
add_run(p, "Hot Deck", bold=True)
add_run(p, " reemplaza cada valor faltante con el valor observado de un ")
add_run(p, '"donante"', bold=True)
add_run(p, ": una observacion similar que si tiene respuesta.")

# Pasos
pasos = [
    ("Paso 1", "Definir variables de\nestratificacion\n(sexo, edad, estado)"),
    ("Paso 2", "Agrupar observaciones\nen celdas segun\nesas variables"),
    ("Paso 3", "Asignar al NA el valor\nde un registro aleatorio\ncon dato"),
]
for i, (titulo, desc) in enumerate(pasos):
    cx = Inches(0.4) + Inches(i * 3.2)
    add_rounded_box(s, cx, Inches(2.0), Inches(2.9), Inches(1.4),
                    BG_HIGHLIGHT)
    add_textbox(s, titulo, cx + Inches(0.1), Inches(2.1),
                Inches(2.7), Inches(0.35),
                font_size=14, bold=True, color=AZUL,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, desc, cx + Inches(0.1), Inches(2.5),
                Inches(2.7), Inches(0.8),
                font_size=12, alignment=PP_ALIGN.CENTER)

# Codigo
add_textbox(s, "Implementacion en R:", Inches(0.5), Inches(3.6),
            Inches(9), Inches(0.3), font_size=14, bold=True)

code_box = add_rounded_box(s, Inches(0.5), Inches(3.95), Inches(9),
                            Inches(2.9), RGBColor(0x1E, 0x2A, 0x3A))
tf_code = add_rich_textbox(s, Inches(0.8), Inches(4.05), Inches(8.4),
                            Inches(2.7))
code_lines_hd = [
    "# --- Metodo 1: Usando el paquete VIM ---",
    "library(VIM)",
    'datos_imp <- hotdeck(',
    '  data = datos_enoe,',
    '  variable = "ingocup",              # Variable a imputar',
    '  ord_var = c("sex", "eda", "cs_p13_1"),  # Ordenamiento',
    '  domain_var = "ent"                 # Dominio (estrato)',
    ")",
    "",
    "# --- Metodo 2: Manual con tidyverse ---",
    "datos_imp <- datos_enoe %>%",
    "  group_by(sex, niv_edu, ent) %>%",
    "  mutate(ingocup_imp = ifelse(",
    "    is.na(ingocup),",
    "    sample(ingocup[!is.na(ingocup)], 1),",
    "    ingocup)) %>% ungroup()",
]
for i, line in enumerate(code_lines_hd):
    p = tf_code.paragraphs[0] if i == 0 else new_paragraph(
        tf_code, space_before=0, space_after=0)
    c = RGBColor(0x7F, 0x8C, 0x8D) if line.startswith("#") else \
        RGBColor(0xE0, 0xE0, 0xE0)
    add_run(p, line, font_size=10, color=c, font_name="Courier New")

# Nota ventaja
add_textbox(s, "Ventaja clave: los valores imputados son valores reales "
            "que alguien reporto, por lo que respetan la distribucion "
            "natural de los datos.",
            Inches(0.5), Inches(7.0), Inches(9), Inches(0.4),
            font_size=12, color=GRIS)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 23: EJERCICIOS survey_mean / survey_total
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_full_banner(s, "Ejercicios: survey_mean y survey_total", y=Inches(0.3))

# Preparacion - codigo compacto
add_textbox(s, "Preparacion: cargar datos y declarar diseno de encuesta.",
            Inches(0.5), Inches(1.05), Inches(9), Inches(0.3),
            font_size=13, bold=True)

code_prep = add_rounded_box(s, Inches(0.5), Inches(1.4), Inches(9),
                             Inches(0.9), RGBColor(0x1E, 0x2A, 0x3A))
tf_p = add_rich_textbox(s, Inches(0.8), Inches(1.45), Inches(8.4),
                         Inches(0.8))
prep_lines = [
    "library(tidyverse); library(srvyr)",
    'sdem <- read_csv("SDEMT125.csv"); coe2 <- read_csv("COE2T125.csv")',
    "enoe <- sdem %>% left_join(coe2)",
    "enoe_dis <- enoe %>% as_survey_design(ids=upm, strata=est_d_tri, "
    "weights=fac_tri, nest=TRUE)",
]
for i, line in enumerate(prep_lines):
    p = tf_p.paragraphs[0] if i == 0 else new_paragraph(
        tf_p, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")

# Ejercicio 1
ej1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.3), Inches(2.5), Inches(4.6), Inches(2.2))
ej1.fill.solid()
ej1.fill.fore_color.rgb = BLANCO
ej1.line.color.rgb = AZUL
ej1.line.width = Pt(2)
add_textbox(s, "Ejercicio 1: survey_total", Inches(0.5), Inches(2.6),
            Inches(4.2), Inches(0.3), font_size=14, bold=True, color=AZUL)
add_textbox(s, "Estima el total de personas\nocupadas en Mexico, por sexo.",
            Inches(0.5), Inches(2.95), Inches(4.2), Inches(0.5),
            font_size=12)
c1 = add_rounded_box(s, Inches(0.5), Inches(3.5), Inches(4.2), Inches(0.85),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c1 = add_rich_textbox(s, Inches(0.7), Inches(3.55), Inches(3.8),
                          Inches(0.75))
c1_lines = [
    "enoe_dis %>%",
    "  filter(clase1 == 1) %>%",
    "  group_by(sex) %>%",
    "  summarise(total = survey_total())",
]
for i, line in enumerate(c1_lines):
    p = tf_c1.paragraphs[0] if i == 0 else new_paragraph(
        tf_c1, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")

# Ejercicio 2
ej2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(5.1), Inches(2.5), Inches(4.6), Inches(2.2))
ej2.fill.solid()
ej2.fill.fore_color.rgb = BLANCO
ej2.line.color.rgb = VERDE
ej2.line.width = Pt(2)
add_textbox(s, "Ejercicio 2: survey_total", Inches(5.3), Inches(2.6),
            Inches(4.2), Inches(0.3), font_size=14, bold=True, color=VERDE)
add_textbox(s, "Estima el total de desocupados\npor entidad federativa.",
            Inches(5.3), Inches(2.95), Inches(4.2), Inches(0.5),
            font_size=12)
c2 = add_rounded_box(s, Inches(5.3), Inches(3.5), Inches(4.2), Inches(0.85),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c2 = add_rich_textbox(s, Inches(5.5), Inches(3.55), Inches(3.8),
                          Inches(0.75))
c2_lines = [
    "enoe_dis %>%",
    "  filter(clase1 == 2) %>%",
    "  group_by(ent) %>%",
    "  summarise(total = survey_total())",
]
for i, line in enumerate(c2_lines):
    p = tf_c2.paragraphs[0] if i == 0 else new_paragraph(
        tf_c2, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")

# Ejercicio 3
ej3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.3), Inches(4.9), Inches(4.6), Inches(2.3))
ej3.fill.solid()
ej3.fill.fore_color.rgb = BLANCO
ej3.line.color.rgb = NARANJA
ej3.line.width = Pt(2)
add_textbox(s, "Ejercicio 3: survey_mean (proporcion)",
            Inches(0.5), Inches(5.0),
            Inches(4.2), Inches(0.3), font_size=14, bold=True, color=NARANJA)
add_textbox(s, "Calcula la tasa de informalidad\nlaboral por entidad.",
            Inches(0.5), Inches(5.35), Inches(4.2), Inches(0.5),
            font_size=12)
c3 = add_rounded_box(s, Inches(0.5), Inches(5.9), Inches(4.2), Inches(1.0),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c3 = add_rich_textbox(s, Inches(0.7), Inches(5.95), Inches(3.8),
                          Inches(0.9))
c3_lines = [
    "enoe_dis %>%",
    "  filter(clase1 == 1) %>%",
    "  mutate(inf = as.numeric(emp_ppal==1)) %>%",
    "  group_by(ent) %>%",
    "  summarise(tasa = survey_mean(inf, na.rm=T))",
]
for i, line in enumerate(c3_lines):
    p = tf_c3.paragraphs[0] if i == 0 else new_paragraph(
        tf_c3, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")

# Ejercicio 4
ej4 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(5.1), Inches(4.9), Inches(4.6), Inches(2.3))
ej4.fill.solid()
ej4.fill.fore_color.rgb = BLANCO
ej4.line.color.rgb = MORADO
ej4.line.width = Pt(2)
add_textbox(s, "Ejercicio 4: survey_mean (media)",
            Inches(5.3), Inches(5.0),
            Inches(4.2), Inches(0.3), font_size=14, bold=True, color=MORADO)
add_textbox(s, "Calcula el ingreso promedio por\nnivel educativo.",
            Inches(5.3), Inches(5.35), Inches(4.2), Inches(0.5),
            font_size=12)
c4 = add_rounded_box(s, Inches(5.3), Inches(5.9), Inches(4.2), Inches(1.0),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c4 = add_rich_textbox(s, Inches(5.5), Inches(5.95), Inches(3.8),
                          Inches(0.9))
c4_lines = [
    "enoe_dis %>%",
    "  filter(clase1 == 1, ingocup > 0) %>%",
    "  group_by(cs_p13_1) %>%",
    "  summarise(ing = survey_mean(",
    "    ingocup, na.rm = TRUE))",
]
for i, line in enumerate(c4_lines):
    p = tf_c4.paragraphs[0] if i == 0 else new_paragraph(
        tf_c4, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 24: EJERCICIOS survey_ratio / survey_quantile
# ══════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_full_banner(s, "Ejercicios: survey_ratio y survey_quantile",
                y=Inches(0.3))

# Ejercicio 5
ej5 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.3), Inches(1.2), Inches(4.6), Inches(2.5))
ej5.fill.solid()
ej5.fill.fore_color.rgb = BLANCO
ej5.line.color.rgb = ROJO
ej5.line.width = Pt(2)
add_textbox(s, "Ejercicio 5: survey_ratio", Inches(0.5), Inches(1.3),
            Inches(4.2), Inches(0.3), font_size=14, bold=True, color=ROJO)
add_textbox(s, "Calcula el ingreso por hora\ntrabajada, por sexo.",
            Inches(0.5), Inches(1.65), Inches(4.2), Inches(0.5),
            font_size=12)
c5 = add_rounded_box(s, Inches(0.5), Inches(2.2), Inches(4.2), Inches(1.2),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c5 = add_rich_textbox(s, Inches(0.7), Inches(2.25), Inches(3.8),
                          Inches(1.1))
c5_lines = [
    "enoe_dis %>%",
    "  filter(clase1==1, ingocup>0, hrsocup>0) %>%",
    "  group_by(sex) %>%",
    "  summarise(ing_hr = survey_ratio(",
    "    ingocup, hrsocup * 4.33, na.rm=T))",
]
for i, line in enumerate(c5_lines):
    p = tf_c5.paragraphs[0] if i == 0 else new_paragraph(
        tf_c5, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")
add_textbox(s, "Existe brecha salarial por hora?",
            Inches(0.5), Inches(3.45), Inches(4.2), Inches(0.25),
            font_size=10, color=GRIS)

# Ejercicio 6
ej6 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(5.1), Inches(1.2), Inches(4.6), Inches(2.5))
ej6.fill.solid()
ej6.fill.fore_color.rgb = BLANCO
ej6.line.color.rgb = RGBColor(0x1A, 0xBC, 0x9C)
ej6.line.width = Pt(2)
add_textbox(s, "Ejercicio 6: survey_ratio", Inches(5.3), Inches(1.3),
            Inches(4.2), Inches(0.3), font_size=14, bold=True,
            color=RGBColor(0x1A, 0xBC, 0x9C))
add_textbox(s, "Calcula la razon informal/formal\npor entidad.",
            Inches(5.3), Inches(1.65), Inches(4.2), Inches(0.5),
            font_size=12)
c6 = add_rounded_box(s, Inches(5.3), Inches(2.2), Inches(4.2), Inches(1.2),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c6 = add_rich_textbox(s, Inches(5.5), Inches(2.25), Inches(3.8),
                          Inches(1.1))
c6_lines = [
    "enoe_dis %>% filter(clase1 == 1) %>%",
    "  mutate(inf = as.numeric(emp_ppal==1),",
    "    form = as.numeric(emp_ppal==2)) %>%",
    "  group_by(ent) %>%",
    "  summarise(r = survey_ratio(inf,form,na.rm=T))",
]
for i, line in enumerate(c6_lines):
    p = tf_c6.paragraphs[0] if i == 0 else new_paragraph(
        tf_c6, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")
add_textbox(s, "Razon > 1 = mas informales que formales.",
            Inches(5.3), Inches(3.45), Inches(4.2), Inches(0.25),
            font_size=10, color=GRIS)

# Ejercicio 7
ej7 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.3), Inches(4.0), Inches(4.6), Inches(2.5))
ej7.fill.solid()
ej7.fill.fore_color.rgb = BLANCO
ej7.line.color.rgb = RGBColor(0x24, 0x71, 0xA3)
ej7.line.width = Pt(2)
add_textbox(s, "Ejercicio 7: survey_quantile", Inches(0.5), Inches(4.1),
            Inches(4.2), Inches(0.3), font_size=14, bold=True,
            color=RGBColor(0x24, 0x71, 0xA3))
add_textbox(s, "Calcula percentiles 25, 50 y 75\ndel ingreso por ocupacion.",
            Inches(0.5), Inches(4.45), Inches(4.2), Inches(0.5),
            font_size=12)
c7 = add_rounded_box(s, Inches(0.5), Inches(5.0), Inches(4.2), Inches(1.1),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c7 = add_rich_textbox(s, Inches(0.7), Inches(5.05), Inches(3.8),
                          Inches(1.0))
c7_lines = [
    "enoe_dis %>%",
    "  filter(clase1 == 1, ingocup > 0) %>%",
    "  summarise(p = survey_quantile(",
    "    ingocup, c(0.25, 0.50, 0.75),",
    "    na.rm = TRUE))",
]
for i, line in enumerate(c7_lines):
    p = tf_c7.paragraphs[0] if i == 0 else new_paragraph(
        tf_c7, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")
add_textbox(s, "Cual es la mediana del ingreso?",
            Inches(0.5), Inches(6.15), Inches(4.2), Inches(0.25),
            font_size=10, color=GRIS)

# Ejercicio 8
ej8 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(5.1), Inches(4.0), Inches(4.6), Inches(2.5))
ej8.fill.solid()
ej8.fill.fore_color.rgb = BLANCO
ej8.line.color.rgb = NARANJA
ej8.line.width = Pt(2)
add_textbox(s, "Ejercicio 8: survey_quantile", Inches(5.3), Inches(4.1),
            Inches(4.2), Inches(0.3), font_size=14, bold=True, color=NARANJA)
add_textbox(s, "Calcula la mediana del ingreso\npor sexo y nivel educativo.",
            Inches(5.3), Inches(4.45), Inches(4.2), Inches(0.5),
            font_size=12)
c8 = add_rounded_box(s, Inches(5.3), Inches(5.0), Inches(4.2), Inches(1.1),
                      RGBColor(0x1E, 0x2A, 0x3A))
tf_c8 = add_rich_textbox(s, Inches(5.5), Inches(5.05), Inches(3.8),
                          Inches(1.0))
c8_lines = [
    "enoe_dis %>%",
    "  filter(clase1 == 1, ingocup > 0) %>%",
    "  group_by(sex, cs_p13_1) %>%",
    "  summarise(med = survey_quantile(",
    "    ingocup, 0.5, na.rm = TRUE))",
]
for i, line in enumerate(c8_lines):
    p = tf_c8.paragraphs[0] if i == 0 else new_paragraph(
        tf_c8, space_before=0, space_after=0)
    add_run(p, line, font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
            font_name="Courier New")
add_textbox(s, "Donde se observa mayor brecha hombre-mujer?",
            Inches(5.3), Inches(6.15), Inches(4.2), Inches(0.25),
            font_size=10, color=GRIS)

# Tip final
box = add_rounded_box(s, Inches(0.5), Inches(6.7), Inches(9), Inches(0.6),
                       BG_HIGHLIGHT)
add_colored_box(s, Inches(0.5), Inches(6.7), Inches(0.08), Inches(0.6), AZUL)
tf_tip = add_rich_textbox(s, Inches(0.8), Inches(6.75), Inches(8.5),
                           Inches(0.5))
p = tf_tip.paragraphs[0]
add_run(p, "Tip: ", font_size=12, bold=True)
add_run(p, "Todas estas funciones devuelven automaticamente el error "
        "estandar. Usa ", font_size=12)
add_run(p, "vartype = \"ci\"", font_size=12, bold=True,
        font_name="Courier New")
add_run(p, " para obtener intervalos de confianza.", font_size=12)


# ══════════════════════════════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════════════════════════════
out_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "Clase07_ENOE.pptx"
)
prs.save(out_path)
print(f"Presentacion guardada en: {out_path}")
print(f"Total de diapositivas: {len(prs.slides)}")
